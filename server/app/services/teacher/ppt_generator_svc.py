import json
import asyncio
from datetime import datetime
from typing import Dict, List, Optional, Any

from fastapi import HTTPException

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.config import SERVER_DIR
from app.core.logger import setup_logger
from app.core.api_client import MODEL_CONFIGS
from app.schemas.teacher.ppt_generator_sch import (
    PPTGenerationRequest, PPTGenerationResponse, PPTSlide,
    PPTOutlineResponse, PPTGenerationFromOutlineRequest, AIModel
)

# 设置日志
logger = setup_logger("ppt_generator_service")

# PPT文件存储目录
PPT_FILES_DIR = SERVER_DIR / "app" / "documents" / "ppt_files"
PPT_OUTLINE_MD_DIR = SERVER_DIR / "app" / "documents" / "ppt_outlines" / "markdown"
PPT_OUTLINE_JSON_DIR = SERVER_DIR / "app" / "documents" / "ppt_outlines" / "json"
PPT_FILES_DIR.mkdir(exist_ok=True, parents=True)
PPT_OUTLINE_MD_DIR.mkdir(exist_ok=True, parents=True)
PPT_OUTLINE_JSON_DIR.mkdir(exist_ok=True, parents=True)



def get_client_and_model(model: AIModel) -> tuple:
    """
    根据模型类型获取对应的客户端和模型名称
    """
    config = MODEL_CONFIGS.get(model)
    if not config:
        logger.warning(f"未知模型类型: {model}, 使用默认 KIMI")
        config = MODEL_CONFIGS[AIModel.KIMI]

    return config["client"], config["model_name"], config["display_name"]



def _build_structured_ppt_prompt(request: PPTGenerationRequest) -> str:
    """构建生成结构化PPT内容的提示词"""
    return f"""你是一位经验丰富的教师和课件专家，请为以下教学内容设计一个详细且内容丰富的PPT结构化数据:

标题: {request.title}
学科: {request.subject}
目标年级: {request.target_grade}
教学目标: {request.teaching_target}
教学重点: {', '.join(request.key_points)}
幻灯片数量: {request.slide_count}
{f"其他信息: {request.additional_info}" if request.additional_info else ""}

请生成{request.slide_count}张幻灯片的结构化数据，每张幻灯片包含以下字段：
- slide_number: 幻灯片序号（1开始）
- slide_type: 幻灯片类型（"cover"=封面、"objective"=学习目标、"content"=内容讲解、"example"=示例代码、"exercise"=练习活动、"summary"=总结）
- title: 幻灯片标题
- main_content: 主要内容（知识点、概念、定义等）
- bullet_points: 要点列表（数组格式）
- code_example: 示例代码（如果有的话，使用具体可执行的代码）
- exercise_content: 练习内容（如果是练习类型幻灯片）
- key_takeaways: 关键要点（数组格式）
- teacher_notes: 教师备注

要求：
1. 内容具体、丰富、实用，避免空洞的表述
2. 对于编程相关内容，提供实际的代码片段
3. 对于概念解释，给出明确的定义和具体例子
4. 确保生成符合要求的幻灯片数量
5. 每个字段都要有内容，如果某字段不适用则填入空字符串或空数组

请以JSON数组格式返回，只返回JSON数据，不要有其他说明文字。

示例格式：
[
  {{
    "slide_number": 1,
    "slide_type": "cover",
    "title": "课程标题",
    "main_content": "课程简介内容",
    "bullet_points": ["要点1", "要点2"],
    "code_example": "",
    "exercise_content": "",
    "key_takeaways": ["关键点1", "关键点2"],
    "teacher_notes": "教师备注内容"
  }}
]"""



def _convert_json_to_markdown(slides_data: List[Dict[str, Any]]) -> str:
    """将结构化JSON数据转换为Markdown格式"""
    md_lines = []

    for slide in slides_data:
        slide_num = slide.get("slide_number", 1)
        slide_type = slide.get("slide_type", "content")
        title = slide.get("title", "")

        # 幻灯片分隔和标题
        md_lines.append(f"# 幻灯片 {slide_num}")
        md_lines.append(f"## {title}")
        md_lines.append("")

        # 主要内容
        main_content = slide.get("main_content", "")
        if main_content:
            md_lines.append("### 主要内容")
            md_lines.append(main_content)
            md_lines.append("")

        # 要点列表
        bullet_points = slide.get("bullet_points", [])
        if bullet_points:
            md_lines.append("### 要点")
            for point in bullet_points:
                md_lines.append(f"- {point}")
            md_lines.append("")

        # 示例代码
        code_example = slide.get("code_example", "")
        if code_example:
            md_lines.append("### 示例代码")
            md_lines.append("```")
            md_lines.append(code_example)
            md_lines.append("```")
            md_lines.append("")

        # 练习内容
        exercise_content = slide.get("exercise_content", "")
        if exercise_content:
            md_lines.append("### 练习")
            md_lines.append(exercise_content)
            md_lines.append("")

        # 关键要点
        key_takeaways = slide.get("key_takeaways", [])
        if key_takeaways:
            md_lines.append("### 关键要点")
            for takeaway in key_takeaways:
                md_lines.append(f"- **{takeaway}**")
            md_lines.append("")

        # 教师备注
        teacher_notes = slide.get("teacher_notes", "")
        if teacher_notes:
            md_lines.append("## 教师备注")
            md_lines.append(teacher_notes)
            md_lines.append("")

        md_lines.append("---")
        md_lines.append("")

    return "\n".join(md_lines)



async def _call_ai_with_retry(client, model_name: str, messages: List[Dict], max_retries: int = 3,
                              timeout: int = 120) -> str:
    """
    带重试机制的AI API调用

    Args:
        client: AI客户端
        model_name: 模型名称
        messages: 消息列表
        max_retries: 最大重试次数
        timeout: 超时时间（秒）

    Returns:
        AI返回的内容

    Raises:
        Exception: 当所有重试都失败时
    """
    last_exception = None

    for attempt in range(max_retries):
        try:
            logger.info(f"正在调用AI接口... (尝试 {attempt + 1}/{max_retries})")

            # 根据幻灯片数量动态调整token数量
            slide_count = 10  # 默认值
            for msg in messages:
                if "幻灯片数量:" in msg.get("content", ""):
                    try:
                        # 从消息中提取幻灯片数量
                        content = msg["content"]
                        start = content.find("幻灯片数量:") + len("幻灯片数量:")
                        end = content.find("\n", start)
                        if end == -1:
                            end = len(content)
                        slide_count = int(content[start:end].strip())
                    except:
                        pass

            # 动态计算max_tokens，每张幻灯片大约需要500-800 tokens
            max_tokens = min(max(slide_count * 600, 4096), 16384)

            response = client.chat.completions.create(
                model=model_name,
                messages=messages,
                temperature=0.6,
                max_tokens=max_tokens,
                stream=False,
                timeout=timeout
            )

            content = response.choices[0].message.content.strip()
            if content:
                logger.info(f"AI接口调用成功 (尝试 {attempt + 1})")
                return content
            else:
                raise Exception("AI返回空内容")

        except Exception as e:
            last_exception = e
            logger.warning(f"第{attempt + 1}次AI调用失败: {str(e)}")

            if attempt < max_retries - 1:
                # 指数退避重试
                wait_time = 2 ** attempt
                logger.info(f"等待{wait_time}秒后重试...")
                await asyncio.sleep(wait_time)
            else:
                logger.error(f"AI调用失败，已达到最大重试次数")

    # 所有重试都失败
    raise Exception(f"AI调用失败，已重试{max_retries}次。最后错误: {str(last_exception)}")



async def generate_ppt_outline(request: PPTGenerationRequest, staff_id: str) -> PPTOutlineResponse:
    """
    生成PPT的结构化JSON数据，然后转换为Markdown格式大纲
    """
    logger.info(f"开始生成PPT大纲: 标题={request.title}, 模型={request.model.value}")

    try:
        # 获取对应模型的客户端和模型名称
        client, model_name, display_name = get_client_and_model(request.model)

        # 构建结构化PPT内容的提示词
        structured_prompt = _build_structured_ppt_prompt(request)

        # 设置消息
        messages = [
            {"role": "system",
             "content": "你是一个专业的教育资源制作助手，擅长生成教学PPT内容。请严格按照要求的JSON格式返回数据。"},
            {"role": "user", "content": structured_prompt}
        ]

        # 根据幻灯片数量调整超时时间
        timeout = max(60, request.slide_count * 10)  # 每张幻灯片至少10秒，最少60秒

        # 调用增强版AI接口
        json_content = await _call_ai_with_retry(
            client=client,
            model_name=model_name,
            messages=messages,
            max_retries=3,
            timeout=timeout
        )

        logger.info(f"成功从{display_name}API获取结构化内容")

        # 解析JSON数据
        try:
            # 移除可能的代码块标记
            if json_content.startswith('```json'):
                json_content = json_content[7:]
            if json_content.endswith('```'):
                json_content = json_content[:-3]

            slides_data = json.loads(json_content)
            logger.info(f"成功解析JSON数据，包含{len(slides_data)}张幻灯片")
        except json.JSONDecodeError as e:
            logger.error(f"JSON解析失败: {str(e)}")
            raise Exception(f"AI返回的数据格式错误: {str(e)}")

        # 将JSON转换为Markdown格式
        md_content = _convert_json_to_markdown(slides_data)
        # 保存JSON和Markdown文件
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # 保存JSON文件
        json_path = PPT_OUTLINE_JSON_DIR / f"outline_json_{staff_id}_{timestamp}_{request.title}.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(slides_data, f, ensure_ascii=False, indent=2)
        logger.info(f"结构化数据已保存到: {json_path}")

        # 保存Markdown文件
        outline_path = PPT_OUTLINE_MD_DIR / f"outline_md_{staff_id}_{timestamp}_{request.title}.md"
        with open(outline_path, "w", encoding="utf-8") as f:
            f.write(md_content)
        logger.info(f"大纲已保存到: {outline_path}")

        return PPTOutlineResponse(
            title=request.title,
            outline_md=md_content,
            model=display_name
        )

    except Exception as e:
        logger.error(f"生成PPT大纲失败: {str(e)}")
        raise Exception(f"生成PPT大纲失败: {str(e)}")



def _try_load_corresponding_json(staff_id: str, title: str) -> Optional[List[Dict[str, Any]]]:
    """尝试加载对应的JSON文件"""
    try:
        # 查找匹配的JSON文件
        pattern = f"outline_json_{staff_id}_*_{title}.json"
        json_files = list(PPT_OUTLINE_JSON_DIR.glob(pattern))

        if json_files:
            # 使用最新的文件
            latest_file = max(json_files, key=lambda f: f.stat().st_mtime)
            with open(latest_file, 'r', encoding='utf-8') as f:
                return json.load(f)
    except Exception as e:
        logger.warning(f"加载JSON文件失败: {str(e)}")

    return None



def _convert_json_to_slide_data(json_data: List[Dict[str, Any]]) -> List[Dict[str, str]]:
    """将JSON结构化数据转换为幻灯片数据"""
    slides_data = []

    for slide_json in json_data:
        slide_data = {
            "title": slide_json.get("title", ""),
            "content": "",
            "note": slide_json.get("teacher_notes", "")
        }

        # 构建内容
        content_parts = []

        # 主要内容
        main_content = slide_json.get("main_content", "")
        if main_content:
            content_parts.append(main_content)

        # 要点列表
        bullet_points = slide_json.get("bullet_points", [])
        if bullet_points:
            content_parts.append("")  # 空行分隔
            for point in bullet_points:
                content_parts.append(f"• {point}")

        # 示例代码
        code_example = slide_json.get("code_example", "")
        if code_example:
            content_parts.append("")
            content_parts.append("**示例代码：**")
            content_parts.append("```")
            content_parts.append(code_example)
            content_parts.append("```")

        # 练习内容
        exercise_content = slide_json.get("exercise_content", "")
        if exercise_content:
            content_parts.append("")
            content_parts.append("**练习：**")
            content_parts.append(exercise_content)

        # 关键要点
        key_takeaways = slide_json.get("key_takeaways", [])
        if key_takeaways:
            content_parts.append("")
            content_parts.append("**关键要点：**")
            for takeaway in key_takeaways:
                content_parts.append(f"★ {takeaway}")

        slide_data["content"] = "\n".join(content_parts)
        slides_data.append(slide_data)

    return slides_data



async def generate_ppt_from_outline(
        request: PPTGenerationFromOutlineRequest,
        staff_id: str
) -> PPTGenerationResponse:
    """从修改后的大纲生成PPT"""
    logger.info(f"从大纲生成PPT: 标题={request.title}, 教师ID={staff_id}")

    try:
        # 尝试找到对应的JSON文件
        json_data = _try_load_corresponding_json(staff_id, request.title)

        if json_data:
            # 使用JSON数据生成PPT
            slides_data = _convert_json_to_slide_data(json_data)
            logger.info("使用JSON结构化数据生成PPT")
        else:
            # 回退到解析Markdown的方式
            slides_data = parse_markdown_outline(request.outline_md)
            logger.info("使用Markdown解析方式生成PPT")

        # 生成文件名
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        file_name = f"ppt_{staff_id}_{timestamp}_{request.title.replace(' ', '_')}.pptx"

        # 创建PPT文件
        file_path = create_pptx_from_structured_data(slides_data, file_name)

        # 构建返回结果
        slides = [
            PPTSlide(
                title=slide["title"],
                content=slide["content"],
                note=slide.get("note", "")
            )
            for slide in slides_data
        ]

        return PPTGenerationResponse(
            title=request.title,
            slides=slides,
            filename=file_name
        )

    except Exception as e:
        logger.error(f"从大纲生成PPT失败: {str(e)}")
        raise Exception(f"从大纲生成PPT失败: {str(e)}")



def create_pptx_from_structured_data(slides_data: List[Dict[str, str]], file_name: str) -> str:
    """
    基于结构化数据创建PPT文件
    """
    try:
        # 创建演示文稿
        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        for slide_data in slides_data:
            # 新建空白幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 标题文本框
            left = Inches(0.5)
            top = Inches(0.4)
            width = slide_width - Inches(1)
            height = Inches(1.2)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_tf = title_box.text_frame
            title_tf.text = slide_data["title"]
            title_p = title_tf.paragraphs[0]
            title_p.font.size = Pt(32)
            title_p.font.bold = True
            title_p.font.name = "微软雅黑"
            title_p.font.color.rgb = RGBColor(31, 56, 100)
            title_p.alignment = PP_ALIGN.LEFT

            # 内容文本框
            content_top = top + height + Inches(0.2)
            content_height = slide_height - content_top - Inches(0.5)
            content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
            content_tf = content_box.text_frame
            content_tf.word_wrap = True

            # 解析并添加内容
            _parse_and_add_content(slide, content_tf, slide_data["content"], left, content_top, width)

            # 添加备注
            if slide_data.get("note"):
                slide.notes_slide.notes_text_frame.text = slide_data["note"]

        # 保存文件
        file_path = PPT_FILES_DIR / file_name
        prs.save(str(file_path))

        logger.info(f"PPT文件已保存: {file_path}")
        return str(file_path)

    except Exception as e:
        logger.error(f"创建PPT失败: {str(e)}")
        raise Exception(f"创建PPT失败: {str(e)}")



def _parse_and_add_content(slide, content_tf, content: str, left, top, width):
    """解析内容并添加到幻灯片"""
    lines = content.split('\n')
    in_code = False
    code_lines = []
    current_y_offset = 0

    for line in lines:
        line = line.strip()

        # 处理代码块
        if line == '```':
            if in_code:
                # 结束代码块
                if code_lines:
                    current_y_offset += _add_code_block_to_slide(
                        slide, code_lines, left, top + Inches(current_y_offset), width
                    )
                code_lines = []
                in_code = False
            else:
                in_code = True
            continue

        if in_code:
            code_lines.append(line)
            continue

        if not line:
            continue

        # 处理不同类型的内容
        if line.startswith('• ') or line.startswith('★ '):
            # 列表项
            p = content_tf.add_paragraph()
            p.text = line[2:]  # 移除符号
            p.level = 0
            p.font.size = Pt(20)
            p.font.name = "微软雅黑"
            if line.startswith('★ '):
                p.font.bold = True
                p.font.color.rgb = RGBColor(200, 50, 50)
        elif line.startswith('**') and line.endswith('**'):
            # 加粗标题
            p = content_tf.add_paragraph()
            p.text = line.replace('**', '')
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = "微软雅黑"
            p.font.color.rgb = RGBColor(50, 100, 150)
        else:
            # 普通文本
            p = content_tf.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.font.name = "微软雅黑"



def _add_code_block_to_slide(slide, code_lines: List[str], left, top, width) -> float:
    """在幻灯片上添加代码块并返回占用的高度（英寸）"""
    if not code_lines:
        return 0

    code_text = "\n".join(code_lines)
    # 估算高度
    height = Pt(18) * (len(code_lines) + 1) + Inches(0.3)

    code_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    code_box.line.color.rgb = RGBColor(200, 200, 200)

    code_tf = code_box.text_frame
    code_tf.word_wrap = True
    code_tf.margin_left = Inches(0.1)
    code_tf.margin_top = Inches(0.1)

    p = code_tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(16)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(60, 60, 60)

    # 返回占用的高度（转换为英寸）
    return height / 914400  # EMU到英寸的转换



# 保持原有的parse_markdown_outline函数作为备用
def parse_markdown_outline(outline_md: str) -> List[Dict[str, str]]:
    """解析Markdown大纲为幻灯片数据列表"""
    slides = []
    current_slide = None
    lines = outline_md.strip().split('\n')

    i = 0
    while i < len(lines):
        line = lines[i].strip()

        # 处理幻灯片标题（以 "# 幻灯片" 开头的行）
        if line.startswith('# 幻灯片') or line.startswith('## 幻灯片'):
            # 保存之前处理的幻灯片
            if current_slide:
                slides.append(current_slide)

            # 创建新的幻灯片数据
            current_slide = {"title": "", "content": "", "note": ""}

            # 查找下一行作为标题（通常是 ## 标题）
            if i + 1 < len(lines) and lines[i + 1].strip().startswith('##'):
                current_slide["title"] = lines[i + 1].strip().replace('##', '').strip()
                i += 2  # 跳过标题行
            else:
                # 如果没有独立的标题行，则使用当前行的内容作为标题
                title_parts = line.split('：', 1)
                if len(title_parts) > 1:
                    current_slide["title"] = title_parts[1].strip()
                else:
                    current_slide["title"] = line.replace('#', '').strip()
                i += 1

            continue

        # 处理分隔符
        if line == '---':
            i += 1
            continue

        # 处理备注
        if line.startswith('## 教师备注'):
            # 进入教师备注部分
            i += 1
            while i < len(lines) and not (lines[i].startswith('# ') or lines[i].startswith('## 幻灯片')):
                note_line = lines[i].strip()
                if note_line and not note_line == '---':
                    if current_slide:
                        current_slide["note"] += note_line + "\n"
                i += 1
            continue

        # 处理普通内容
        if current_slide is not None and line:
            current_slide["content"] += line + "\n"

        i += 1

    # 添加最后一个幻灯片
    if current_slide:
        slides.append(current_slide)

    logger.info(f"成功解析出 {len(slides)} 张幻灯片")
    return slides



async def list_ppt_outlines_service(staff_id: str):
    """
    列出教师的所有PPT大纲

    Args:
        staff_id: 教师工号

    Returns:
        包含大纲列表的字典
    """
    outlines = []

    for file_path in PPT_OUTLINE_MD_DIR.glob(f"outline_md_{staff_id}_*.md"):
        try:
            filename = file_path.name
            base_name = filename.rsplit('.', 1)[0]
            parts = base_name.split('_')

            # 修正文件名解析逻辑
            # 文件名格式: outline_md_{staff_id}_{timestamp}_{title}.md
            if len(parts) >= 4:
                # parts[0] = "outline", parts[1] = "md", parts[2] = staff_id, parts[3] = timestamp
                title = '_'.join(parts[4:]) if len(parts) > 4 else "未知标题"
            else:
                title = '未知标题'

            # 读取文件内容
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()

            outlines.append({
                "filename": filename,
                "title": title,
                "created_at": file_path.stat().st_ctime,
                "preview": content[:50] + ("..." if len(content) > 50 else ""),
            })
        except Exception as e:
            logger.error(f"处理大纲文件 {file_path.name} 失败: {str(e)}")
            continue

    return {"outlines": outlines}



async def download_ppt_outline_service(file_name: str, staff_id: str):
    """
    处理PPT大纲下载的权限验证和文件检查

    Args:
        file_name: 要下载的大纲文件名
        staff_id: 教师工号

    Returns:
        文件路径对象

    Raises:
        HTTPException: 当权限验证失败或文件不存在时
    """
    # 验证文件权限
    if not file_name.startswith(f"outline_md_{staff_id}_"):
        logger.warning(f"教师(工号:{staff_id}) 尝试访问非自己的大纲文件: {file_name}")
        raise HTTPException(status_code=403, detail="没有权限访问此文件")

    file_path = PPT_OUTLINE_MD_DIR / file_name
    if not file_path.exists():
        logger.warning(f"教师(工号:{staff_id}) 请求的大纲文件不存在: {file_name}")
        raise HTTPException(status_code=404, detail="文件不存在")

    return file_path



async def list_ppt_files_service(staff_id: str):
    """
    列出教师的所有PPT文件

    Args:
        staff_id: 教师工号

    Returns:
        包含文件列表的字典
    """
    staff_id_pattern = f"ppt_{staff_id}_"
    files = []

    for file_path in PPT_FILES_DIR.glob("*.pptx"):
        # 只返回当前教师的文件
        if file_path.name.startswith(staff_id_pattern):
            files.append({
                "file_name": file_path.name,
                "size": file_path.stat().st_size,
                "created_at": file_path.stat().st_ctime
            })

    return {"files": files}



async def download_ppt_service(file_name: str, staff_id: str):
    """
    处理PPT下载的权限验证和文件检查

    Args:
        file_name: 要下载的文件名
        staff_id: 教师工号

    Returns:
        文件路径对象

    Raises:
        HTTPException: 当权限验证失败或文件不存在时
    """
    # 验证文件权限
    if not file_name.startswith(f"ppt_{staff_id}_"):
        logger.warning(f"教师(工号:{staff_id}) 尝试访问非自己的文件: {file_name}")
        raise HTTPException(status_code=403, detail="没有权限访问此文件")

    file_path = PPT_FILES_DIR / file_name
    if not file_path.exists():
        logger.warning(f"教师(工号:{staff_id}) 请求的文件不存在: {file_name}")
        raise HTTPException(status_code=404, detail="文件不存在")

    return file_path



async def delete_ppt_outline_service(file_name: str, staff_id: str) -> None:
    """
    删除PPT大纲文件，需要同时删除md与json格式的大纲文件

    Args:
        file_name: 要删除的大纲文件名
        staff_id: 教师工号

    Raises:
        HTTPException: 当权限验证失败或文件不存在时
    """
    # 验证文件权限 (格式应为 outline_md_staff_id_*.md)
    if not file_name.startswith(f"outline_md_{staff_id}_"):
        logger.warning(f"权限拒绝: 教师(工号:{staff_id})尝试删除非本人大纲: {file_name}")
        raise HTTPException(status_code=403, detail="您无权删除此大纲")

    # 构建MD文件路径
    md_file_path = PPT_OUTLINE_MD_DIR / file_name
    # 根据MD文件名构建对应的JSON文件名
    # outline_md_staff_id_timestamp_title.md -> outline_json_staff_id_timestamp_title.json
    if file_name.endswith(".md"):
        json_file_name = file_name.replace("outline_md_", "outline_json_").replace(".md", ".json")
        json_file_path = PPT_OUTLINE_JSON_DIR / json_file_name
    else:
        raise HTTPException(status_code=400, detail="不支持的文件类型")

    # 检查MD文件是否存在
    if not md_file_path.exists():
        logger.warning(f"删除大纲失败: 文件不存在 '{md_file_path}'")
        raise HTTPException(status_code=404, detail="大纲文件不存在")

    try:
        # 删除MD文件
        md_file_path.unlink()
        logger.info(f"MD文件已删除: {md_file_path}")

        # 删除对应的JSON文件（如果存在）
        if json_file_path.exists():
            json_file_path.unlink()
            logger.info(f"关联JSON文件已删除: {json_file_path}")
        else:
            logger.info(f"JSON文件不存在，跳过删除: {json_file_path}")

        logger.info(f"教师(工号:{staff_id})成功删除大纲文件: {file_name}")

    except Exception as e:
        logger.error(f"删除大纲文件失败: {str(e)}")
        raise Exception(f"删除大纲文件失败: {str(e)}")



async def delete_ppt_file_service(file_name: str, staff_id: str) -> None:
    """
    删除PPT文件

    Args:
        file_name: 要删除的PPT文件名
        staff_id: 教师工号

    Raises:
        HTTPException: 当权限验证失败或文件不存在时
    """
    # 验证文件权限 (格式应为 ppt_staff_id_*.pptx)
    if not file_name.startswith(f"ppt_{staff_id}_"):
        logger.warning(f"权限拒绝: 教师(工号:{staff_id})尝试删除非本人PPT文件: {file_name}")
        raise HTTPException(status_code=403, detail="您无权删除此文件")

    file_path = PPT_FILES_DIR / file_name
    if not file_path.exists():
        logger.warning(f"删除PPT文件失败: 文件不存在 '{file_path}'")
        raise HTTPException(status_code=404, detail="PPT文件不存在")

    try:
        file_path.unlink()
        logger.info(f"教师(工号:{staff_id})成功删除PPT文件: {file_name}")
    except Exception as e:
        logger.error(f"删除PPT文件失败: {str(e)}")
        raise Exception(f"删除PPT文件失败: {str(e)}")